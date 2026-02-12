"""Vision model service for image analysis and document processing using Groq."""

import base64
import logging
from io import BytesIO
from typing import Optional

from PIL import Image

from app.config import get_settings
from app.llm_client import get_llm_client

# Optional document processing imports
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)
settings = get_settings()


class VisionService:
    """Service for analyzing images with Groq Vision model."""
    
    def __init__(self):
        """Initialize vision service."""
        self.llm_client = get_llm_client()
        self.max_image_size = (1024, 1024)  # Max dimensions
        self.max_file_size = 5 * 1024 * 1024  # 5MB
        
    def validate_image(self, file_bytes: bytes, filename: str) -> tuple[bool, Optional[str]]:
        """
        Validate image file.
        
        Args:
            file_bytes: Image file bytes
            filename: Original filename
            
        Returns:
            (is_valid, error_message)
        """
        # Check file size
        if len(file_bytes) > self.max_file_size:
            return False, f"Image too large. Maximum size is {self.max_file_size // (1024*1024)}MB"
        
        # Check if it's a valid image
        try:
            img = Image.open(BytesIO(file_bytes))
            img.verify()
            
            # Check format
            valid_formats = {'PNG', 'JPEG', 'JPG', 'WEBP', 'GIF'}
            if img.format not in valid_formats:
                return False, f"Unsupported format: {img.format}. Use PNG, JPEG, WEBP, or GIF"
            
            return True, None
            
        except Exception as e:
            logger.error(f"Image validation failed: {e}")
            return False, f"Invalid image file: {str(e)}"
    
    def process_image(self, file_bytes: bytes) -> str:
        """
        Process and optimize image for vision model.
        
        Args:
            file_bytes: Raw image bytes
            
        Returns:
            Base64 encoded image string
        """
        try:
            # Open image
            img = Image.open(BytesIO(file_bytes))
            
            # Convert to RGB if needed (e.g., PNG with transparency)
            if img.mode in ('RGBA', 'LA', 'P'):
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'P':
                    img = img.convert('RGBA')
                background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                img = background
            elif img.mode != 'RGB':
                img = img.convert('RGB')
            
            # Resize if too large
            if img.size[0] > self.max_image_size[0] or img.size[1] > self.max_image_size[1]:
                img.thumbnail(self.max_image_size, Image.Resampling.LANCZOS)
                logger.info(f"Resized image to {img.size}")
            
            # Convert to JPEG for optimal size
            buffer = BytesIO()
            img.save(buffer, format='JPEG', quality=85, optimize=True)
            img_bytes = buffer.getvalue()
            
            # Encode to base64
            base64_image = base64.b64encode(img_bytes).decode('utf-8')
            
            logger.info(f"Processed image: {len(base64_image)} chars base64")
            return base64_image
            
        except Exception as e:
            logger.error(f"Image processing failed: {e}", exc_info=True)
            raise ValueError(f"Failed to process image: {str(e)}")
    
    async def analyze_image(
        self,
        image_base64: str,
        prompt: str = "Describe this image in detail.",
        user_id: str = None
    ) -> dict:
        """
        Analyze image using Groq Vision model.
        
        Args:
            image_base64: Base64 encoded image
            prompt: Analysis prompt
            user_id: Optional user ID for logging
            
        Returns:
            Analysis result with text and metadata
        """
        try:
            logger.info(f"Analyzing image for user {user_id} with prompt: {prompt[:50]}...")
            
            # Prepare messages for vision model
            messages = [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": prompt
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{image_base64}"
                            }
                        }
                    ]
                }
            ]
            
            # Call Groq Vision model (Llama 4 Scout)
            response = await self.llm_client.generate_completion_async(
                messages=messages,
                model="meta-llama/llama-4-scout-17b-16e-instruct",
                temperature=0.7,
                max_tokens=1000
            )
            
            analysis_text = response.strip()
            
            logger.info(f"✅ Image analysis complete: {len(analysis_text)} chars")
            
            return {
                "analysis": analysis_text,
                "model": "meta-llama/llama-4-scout-17b-16e-instruct",
                "prompt": prompt,
                "success": True
            }
            
        except Exception as e:
            logger.error(f"Image analysis failed: {e}", exc_info=True)
            return {
                "analysis": None,
                "model": "meta-llama/llama-4-scout-17b-16e-instruct",
                "prompt": prompt,
                "success": False,
                "error": str(e)
            }
    
    def extract_text_from_pdf(self, file_bytes: bytes) -> str:
        """Extract text from PDF file."""
        if not PDF_AVAILABLE:
            raise ValueError("PDF processing not available. Install PyPDF2: pip install PyPDF2")
        
        try:
            pdf_reader = PyPDF2.PdfReader(BytesIO(file_bytes))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n\n"
            return text.strip()
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            raise ValueError(f"Failed to extract text from PDF: {str(e)}")
    
    def extract_text_from_docx(self, file_bytes: bytes) -> str:
        """Extract text from DOCX file."""
        if not DOCX_AVAILABLE:
            raise ValueError("DOCX processing not available. Install python-docx: pip install python-docx")
        
        try:
            doc = docx.Document(BytesIO(file_bytes))
            text = "\n\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            return text.strip()
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            raise ValueError(f"Failed to extract text from DOCX: {str(e)}")
    
    def extract_text_from_pptx(self, file_bytes: bytes) -> str:
        """Extract text from PPTX file."""
        if not PPTX_AVAILABLE:
            raise ValueError("PPTX processing not available. Install python-pptx: pip install python-pptx")
        
        try:
            prs = Presentation(BytesIO(file_bytes))
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"--- Slide {slide_num} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                text += "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            raise ValueError(f"Failed to extract text from PPTX: {str(e)}")
    
    async def analyze_document(
        self,
        file_bytes: bytes,
        filename: str,
        prompt: str = "Summarize this document.",
        user_id: str = None
    ) -> dict:
        """
        Extract text from document and analyze with LLM.
        
        Args:
            file_bytes: Document file bytes
            filename: Original filename
            prompt: Analysis prompt
            user_id: Optional user ID for logging
            
        Returns:
            Analysis result with text and metadata
        """
        try:
            # Determine file type and extract text
            file_ext = filename.lower().split('.')[-1]
            
            if file_ext == 'pdf':
                text = self.extract_text_from_pdf(file_bytes)
            elif file_ext in ['docx', 'doc']:
                text = self.extract_text_from_docx(file_bytes)
            elif file_ext in ['pptx', 'ppt']:
                text = self.extract_text_from_pptx(file_bytes)
            else:
                raise ValueError(f"Unsupported document type: {file_ext}")
            
            if not text.strip():
                raise ValueError("No text extracted from document")
            
            logger.info(f"Extracted {len(text)} chars from {filename}")
            
            # Analyze with LLM
            full_prompt = f"{prompt}\n\nDocument content:\n{text[:4000]}"  # Limit to 4000 chars
            
            response = await self.llm_client.generate_completion_async(
                messages=[{"role": "user", "content": full_prompt}],
                model=settings.LLM_MODEL,
                temperature=0.7,
                max_tokens=1500
            )
            
            analysis_text = response.strip()
            
            logger.info(f"✅ Document analysis complete: {len(analysis_text)} chars")
            
            return {
                "analysis": analysis_text,
                "extracted_text": text[:500] + "..." if len(text) > 500 else text,
                "model": settings.LLM_MODEL,
                "prompt": prompt,
                "success": True
            }
            
        except Exception as e:
            logger.error(f"Document analysis failed: {e}", exc_info=True)
            return {
                "analysis": None,
                "model": settings.LLM_MODEL,
                "prompt": prompt,
                "success": False,
                "error": str(e)
            }


# Global instance
vision_service = VisionService()
